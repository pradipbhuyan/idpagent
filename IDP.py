# ==============================
# INTELLIGENT DOCUMENT PROCESSOR - CLEAN AGENTIC UI
# IDP.py
# ==============================

import re
import json
import time
import base64
import tempfile
import hashlib
import traceback
from pathlib import Path

import pandas as pd
import streamlit as st

from openai import OpenAI
from docx import Document as DocxDocument
from pptx import Presentation

from langchain_core.documents import Document
from langchain_core.messages import HumanMessage
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from langchain_community.document_loaders import TextLoader, PyPDFLoader

from workflow import build_graph
from core import (
    build_resume,
    json_to_kv_dataframe,
    generate_excel,
    send_to_concur,
    validate_document_data,
    build_confidence_map,
)

# ------------------------------
# PAGE CONFIG
# ------------------------------
st.set_page_config("IDP - Professional", layout="wide")
USERS = st.secrets.get("users", {})

# ------------------------------
# CACHED MODELS
# ------------------------------
@st.cache_resource
def get_llm(api_key, model):
    return ChatOpenAI(model=model, temperature=0, api_key=api_key)

@st.cache_resource
def get_embeddings(api_key):
    return OpenAIEmbeddings(api_key=api_key)

# ------------------------------
# AUTH
# ------------------------------
def validate_api_key(api_key):
    try:
        client = OpenAI(api_key=api_key)
        client.models.list()
        return True
    except Exception:
        return False


def login():
    logo_path = Path(__file__).parent / "IDP-Logo1.png"
    col1, col2, col3 = st.columns([1, 2, 1])

    with col2:
        if logo_path.exists():
            st.image(logo_path, width=220)

        st.markdown("### Sign In")
        username = st.text_input("Username")
        password = st.text_input("Password", type="password")
        api_key = st.text_input("OpenAI API Key", type="password")

        if st.button("Login", use_container_width=True):
            if username not in USERS or USERS[username]["password"] != password:
                st.error("Invalid username or password")
                return

            if not api_key:
                st.error("Please enter your OpenAI API key")
                return

            with st.spinner("Validating API key..."):
                if not validate_api_key(api_key):
                    st.error("Invalid OpenAI API key")
                    return

            st.session_state["logged_in"] = True
            st.session_state["user"] = username
            st.session_state["role"] = USERS[username].get("role", "user")
            st.session_state["api_key"] = api_key
            st.success(f"Welcome {username}")
            st.rerun()

# ------------------------------
# SESSION INIT
# ------------------------------
DEFAULT_METRICS = {
    "tokens": 0,
    "input_tokens": 0,
    "output_tokens": 0,
    "cost": 0.0,
    "response_times": [],
    "calls": 0,
}

DEFAULT_KEYS = {
    "logged_in": False,
    "user": None,
    "role": None,
    "api_key": None,
    "model_choice": "gpt-4o-mini",
    "structured_data": None,
    "doc_type": None,
    "vectorstore": None,
    "full_text": None,
    "generated_resume": None,
    "chat_history": [],
    "suggested_questions": [],
    "metrics": DEFAULT_METRICS.copy(),
    "doc_costs": {},
    "auto_result": None,
    "file_hash": None,
    "current_file": None,
    "processing_error": None,
    "elapsed_time": 0.0,
    "run_started_at": None,
    "agent_events": [],
    "agent_logs": [],
    "agent_status": "Idle",
    "current_step": "Waiting for upload",
    "progress_value": 0,
    "review_data": None,
    "confidence_map": None,
    "validation_result": None,
    "live_step_placeholder": None,
    "live_progress_placeholder": None,
    "live_event_placeholder": None,
}

for k, v in DEFAULT_KEYS.items():
    if k not in st.session_state:
        st.session_state[k] = v

# ------------------------------
# LOGIN GATE
# ------------------------------
if not st.session_state["logged_in"]:
    login()
    st.stop()

if not st.session_state.get("api_key"):
    st.error("API key missing. Please login again.")
    st.stop()

# ------------------------------
# SESSION HELPERS
# ------------------------------
def reset_document_state():
    st.session_state["structured_data"] = None
    st.session_state["doc_type"] = None
    st.session_state["vectorstore"] = None
    st.session_state["full_text"] = None
    st.session_state["generated_resume"] = None
    st.session_state["chat_history"] = []
    st.session_state["suggested_questions"] = []
    st.session_state["auto_result"] = None
    st.session_state["processing_error"] = None
    st.session_state["elapsed_time"] = 0.0
    st.session_state["run_started_at"] = None
    st.session_state["agent_events"] = []
    st.session_state["agent_logs"] = []
    st.session_state["agent_status"] = "Idle"
    st.session_state["current_step"] = "Waiting for upload"
    st.session_state["progress_value"] = 0
    st.session_state["review_data"] = None
    st.session_state["confidence_map"] = None
    st.session_state["validation_result"] = None
    st.session_state["live_step_placeholder"] = None
    st.session_state["live_progress_placeholder"] = None
    st.session_state["live_event_placeholder"] = None


def get_metrics_snapshot():
    m = st.session_state.get("metrics", {})
    return {
        "tokens": m.get("tokens", 0),
        "input_tokens": m.get("input_tokens", 0),
        "output_tokens": m.get("output_tokens", 0),
        "cost": m.get("cost", 0.0),
        "calls": m.get("calls", 0),
    }


def diff_metrics(before, after):
    return {
        "tokens": after.get("tokens", 0) - before.get("tokens", 0),
        "input_tokens": after.get("input_tokens", 0) - before.get("input_tokens", 0),
        "output_tokens": after.get("output_tokens", 0) - before.get("output_tokens", 0),
        "cost": after.get("cost", 0.0) - before.get("cost", 0.0),
        "calls": after.get("calls", 0) - before.get("calls", 0),
    }


def tracked_llm_call(prompt):
    llm = get_llm(
        st.session_state["api_key"],
        st.session_state.get("model_choice", "gpt-4o-mini")
    )

    start = time.time()
    response = llm.invoke(prompt)
    duration = time.time() - start

    try:
        usage = getattr(response, "response_metadata", {}).get("token_usage", {})
        input_tokens = usage.get("prompt_tokens", 0)
        output_tokens = usage.get("completion_tokens", 0)
    except Exception:
        input_tokens = len(str(prompt)) // 4
        output_tokens = len(str(getattr(response, "content", ""))) // 4

    total_tokens = input_tokens + output_tokens
    input_cost = input_tokens * 0.00015 / 1000
    output_cost = output_tokens * 0.0006 / 1000
    total_cost = input_cost + output_cost

    m = st.session_state.metrics
    m["tokens"] += total_tokens
    m["input_tokens"] += input_tokens
    m["output_tokens"] += output_tokens
    m["cost"] += total_cost
    m["calls"] += 1
    m["response_times"].append(duration)

    doc = st.session_state.get("current_file") or "unknown"
    if doc not in st.session_state.doc_costs:
        st.session_state.doc_costs[doc] = {"cost": 0.0, "tokens": 0}

    st.session_state.doc_costs[doc]["cost"] += total_cost
    st.session_state.doc_costs[doc]["tokens"] += total_tokens

    return response


def push_agent_log(message):
    st.session_state.agent_logs.append(message)


def record_agent_event(step, status, message="", duration=None, metrics=None):
    st.session_state.agent_events.append({
        "step": step,
        "status": status,
        "message": message,
        "duration": duration,
        "metrics": metrics or {}
    })
    refresh_live_activity()


def init_live_progress_placeholders():
    if "live_step_placeholder" not in st.session_state:
        st.session_state["live_step_placeholder"] = None
    if "live_progress_placeholder" not in st.session_state:
        st.session_state["live_progress_placeholder"] = None
    if "live_event_placeholder" not in st.session_state:
        st.session_state["live_event_placeholder"] = None


def render_live_activity_snapshot():
    events = st.session_state.get("agent_events", [])
    lines = []

    for event in events[-10:]:
        status = event.get("status", "pending")
        if status == "done":
            icon = "✅"
        elif status == "error":
            icon = "❌"
        elif status == "running":
            icon = "🔄"
        else:
            icon = "⏳"

        line = f"{icon} {event.get('step', '')}"
        if event.get("message"):
            line += f" — {event.get('message')}"
        lines.append(line)

    if not lines:
        lines = ["⏳ Waiting for upload"]

    return "\n\n".join(lines)


def refresh_live_activity():
    step_placeholder = st.session_state.get("live_step_placeholder")
    progress_placeholder = st.session_state.get("live_progress_placeholder")
    event_placeholder = st.session_state.get("live_event_placeholder")

    current_step = st.session_state.get("current_step", "Waiting for upload")
    progress_value = int(st.session_state.get("progress_value", 0))
    events = st.session_state.get("agent_events", [])
    logs = st.session_state.get("agent_logs", [])

    has_started = len(events) > 0 or progress_value > 0 or current_step != "Waiting for upload"

    if step_placeholder is not None:
        if has_started:
            step_placeholder.markdown(f"#### Progress\n\n**Current Step:** {current_step}")
        else:
            step_placeholder.empty()

    if progress_placeholder is not None:
        if has_started:
            progress_placeholder.progress(progress_value)
        else:
            progress_placeholder.empty()

    if event_placeholder is not None:
        if not has_started:
            event_placeholder.empty()
            return

        content = ["#### Completed Steps"]

        real_events = [
            e for e in events
            if e.get("step") and e.get("step").strip().lower() != "waiting for upload"
        ]

        for event in real_events[-8:]:
            status = event.get("status", "pending")
            if status == "done":
                icon = "✅"
            elif status == "error":
                icon = "❌"
            elif status == "running":
                icon = "🔄"
            else:
                icon = "⏳"

            line = f"{icon} **{event.get('step', '')}**"
            if event.get("message"):
                line += f"  \n{event.get('message')}"
            content.append(line)

        if logs:
            content.append("---")
            content.append("**Recent Logs**")
            for log in logs[-5:]:
                content.append(f"- {log}")

        event_placeholder.markdown("\n\n".join(content))
        
def update_progress(percent, message):
    st.session_state["progress_value"] = percent
    st.session_state["current_step"] = message
    st.session_state["agent_status"] = "Running"

    events = st.session_state.get("agent_events", [])
    if not events or events[-1].get("step") != message or events[-1].get("status") != "running":
        events.append({
            "step": message,
            "status": "running",
            "message": "",
            "duration": None,
            "metrics": {}
        })
        st.session_state["agent_events"] = events

    refresh_live_activity()


def save_temp_file(uploaded_file):
    suffix = Path(uploaded_file.name).suffix
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(uploaded_file.getvalue())
        return tmp.name


def load_default_resume_template_bytes():
    possible_paths = [
        Path("templates/resume_template.docx"),
        Path("templates:resume_template.docx"),
        Path(__file__).parent / "templates" / "resume_template.docx",
        Path(__file__).parent / "templates:resume_template.docx",
    ]
    for p in possible_paths:
        if p.exists():
            with open(p, "rb") as f:
                return f.read()
    return None


def extract_docx_text(file_path):
    doc = DocxDocument(file_path)
    parts = []

    for p in doc.paragraphs:
        if p.text and p.text.strip():
            parts.append(p.text.strip())

    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    for section in doc.sections:
        for p in section.header.paragraphs:
            if p.text and p.text.strip():
                parts.append(p.text.strip())
        for p in section.footer.paragraphs:
            if p.text and p.text.strip():
                parts.append(p.text.strip())

    return "\n".join(parts).strip()

# ------------------------------
# FILE PROCESSING
# ------------------------------
def process_file(uploaded_file):
    documents = []
    if not uploaded_file:
        return documents

    uploaded_file.seek(0)
    suffix = Path(uploaded_file.name).suffix.lower()

    if suffix in [".png", ".jpg", ".jpeg"]:
        encoded = base64.b64encode(uploaded_file.getvalue()).decode()
        message = HumanMessage(
            content=[
                {
                    "type": "text",
                    "text": """You are an OCR assistant.

Extract ALL visible text from the image.

Rules:
- Do NOT skip anything
- Preserve numbers, amounts, dates
- Preserve line structure
- Output plain text only
"""
                },
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:image/png;base64,{encoded}"}
                }
            ]
        )

        try:
            llm = get_llm(
                st.session_state["api_key"],
                st.session_state.get("model_choice", "gpt-4o-mini")
            )
            response = llm.invoke([message])
            content = getattr(response, "content", "") or ""
            if not str(content).strip():
                content = "No readable text found in image"
            documents.append(Document(page_content=str(content)))
        except Exception as e:
            st.error(f"OCR failed: {str(e)}")
            return []
        return documents

    file_path = save_temp_file(uploaded_file)

    try:
        if suffix == ".txt":
            try:
                documents.extend(TextLoader(file_path, encoding="utf-8").load())
            except Exception:
                documents.extend(TextLoader(file_path, encoding="cp1252").load())

        elif suffix == ".pdf":
            docs = PyPDFLoader(file_path).load()
            if docs:
                documents.extend(docs)

        elif suffix == ".docx":
            text = extract_docx_text(file_path)
            if text.strip():
                documents.append(Document(page_content=text))

        elif suffix == ".pptx":
            prs = Presentation(file_path)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text and shape.text.strip():
                        text_parts.append(shape.text.strip())
            text = "\n".join(text_parts).strip()
            if text:
                documents.append(Document(page_content=text))

        elif suffix == ".xlsx":
            excel_file = pd.ExcelFile(file_path)
            sheet_texts = []
            for sheet in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet)
                sheet_texts.append(f"Sheet: {sheet}")
                sheet_texts.append(df.to_string(index=False))
            text = "\n\n".join(sheet_texts).strip()
            if text:
                documents.append(Document(page_content=text))

    except Exception as e:
        st.error(f"File processing failed: {str(e)}")
        return []

    return documents


def create_vectorstore(docs):
    if not docs:
        return None

    splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=150)
    chunks = splitter.split_documents(docs)
    if not chunks:
        return None

    for chunk in chunks:
        chunk.metadata = {"source": st.session_state.get("current_file", "unknown")}

    try:
        emb = get_embeddings(st.session_state["api_key"])
        return Chroma.from_documents(chunks, embedding=emb)
    except Exception:
        st.error("Vectorstore creation failed")
        st.code(traceback.format_exc())
        return None


def get_suggested_questions(doc_type):
    if doc_type == "invoice":
        return [
            "What is the total amount?",
            "Who is the vendor?",
            "What is the invoice date?"
        ]
    elif doc_type == "resume":
        return [
            "Summarize this candidate",
            "What skills does the candidate have?",
            "What is the experience?"
        ]
    elif doc_type == "ticket":
        return [
            "What is the ticket number?",
            "What is the travel date?",
            "What are the key details?"
        ]
    return [
        "What is this document?",
        "What are the key points?"
    ]


def normalize_graph_result(result):
    if not isinstance(result, dict):
        return {
            "doc_type": None,
            "structured_data": None,
            "result": {},
            "error": "Graph returned non-dict output"
        }

    doc_type = result.get("doc_type") or result.get("type")
    structured_data = result.get("data") if doc_type in ["invoice", "ticket"] else None
    inner = result.get("result", {}) if isinstance(result.get("result", {}), dict) else {}

    if not inner:
        inner = {
            "type": result.get("result_type") or doc_type,
            "file": result.get("file"),
            "excel": result.get("excel"),
            "table": result.get("table"),
            "payload": result.get("payload"),
            "message": result.get("message"),
            "file_name": result.get("file_name"),
            "data": result.get("data"),
            "concur_status": result.get("concur_status"),
            "concur_mode": result.get("concur_mode"),
        }

    return {
        "doc_type": doc_type,
        "structured_data": structured_data,
        "result": inner,
        "error": result.get("error"),
        "step_metrics": result.get("step_metrics", []),
        "confidence": result.get("confidence"),
        "validation": result.get("validation"),
    }


def process_uploaded_document(uploaded_file):
    current_file = uploaded_file.name
    st.session_state.current_file = current_file
    st.session_state.run_started_at = time.time()
    st.session_state.agent_status = "Running"

    push_agent_log(f"Upload received: {current_file}")
    record_agent_event("Upload received", "done", f"File: {current_file}")
    update_progress(5, "Upload received")

    docs = process_file(uploaded_file)
    if not docs:
        raise ValueError("Failed to process document")

    record_agent_event("Text extraction", "done", "Document converted into searchable text")
    update_progress(15, "Text extracted")

    full_text = "\n".join(
        [str(d.page_content) for d in docs if d is not None and getattr(d, "page_content", None)]
    ).strip()

    if not full_text:
        raise ValueError("No text extracted from document")

    update_progress(20, "Creating search index")
    vectorstore = create_vectorstore(docs)
    record_agent_event("Search index", "done", "Vector index created")
    update_progress(25, "Search index ready")

    graph = build_graph()
    graph_input = {
        "text": full_text,
        "filename": uploaded_file.name,
        "template": load_default_resume_template_bytes(),
        "progress": update_progress,
    }

    before_metrics = get_metrics_snapshot()
    raw_result = graph.invoke(graph_input)
    after_metrics = get_metrics_snapshot()
    auto_metrics = diff_metrics(before_metrics, after_metrics)

    normalized = normalize_graph_result(raw_result)
    doc_type = normalized.get("doc_type")
    structured_data = normalized.get("structured_data")
    result = normalized.get("result", {})
    error = normalized.get("error")
    step_metrics = normalized.get("step_metrics", [])

    if error:
        raise RuntimeError(str(error))

    if not doc_type:
        raise ValueError("Document type could not be determined")

    review_data = result.get("data") or structured_data or {}
    validation = normalized.get("validation") or validate_document_data(review_data, doc_type)
    confidence = normalized.get("confidence") or build_confidence_map(review_data, doc_type)

    st.session_state.full_text = full_text
    st.session_state.vectorstore = vectorstore
    st.session_state.doc_type = doc_type
    st.session_state.structured_data = structured_data if doc_type in ["invoice", "ticket"] else None
    st.session_state.auto_result = {
        "doc_type": doc_type,
        "structured_data": structured_data if doc_type in ["invoice", "ticket"] else None,
        "result": result,
        "metrics": auto_metrics,
        "step_metrics": step_metrics,
    }
    st.session_state.review_data = review_data
    st.session_state.validation_result = validation
    st.session_state.confidence_map = confidence
    st.session_state.suggested_questions = get_suggested_questions(doc_type)

    if doc_type == "resume":
        st.session_state.generated_resume = result.get("file")

    st.session_state.elapsed_time = time.time() - st.session_state.run_started_at
    st.session_state.agent_status = "Done"
    st.session_state.current_step = "Completed"
    st.session_state.progress_value = 100
    refresh_live_activity()

    push_agent_log(f"Processing complete: {doc_type}")
    record_agent_event("Output ready", "done", f"Detected type: {doc_type}")
    update_progress(100, "Completed")

# ------------------------------
# REVIEW HELPERS
# ------------------------------
def refresh_review_scores():
    data = st.session_state.get("review_data") or {}
    doc_type = st.session_state.get("doc_type") or "other"
    st.session_state.validation_result = validate_document_data(data, doc_type)
    st.session_state.confidence_map = build_confidence_map(data, doc_type)


def compact_field(label, value):
    st.markdown(
        f"**{label}**  \n<small>{value if value not in [None, ''] else '-'}</small>",
        unsafe_allow_html=True
    )


def render_validation_summary():
    validation = st.session_state.get("validation_result") or {}
    issues = validation.get("issues", [])
    warnings = validation.get("warnings", [])
    passed = validation.get("passed", False)

    st.markdown("#### Validation")
    if passed:
        st.success("Ready for approval")
    else:
        st.warning("Needs review before approval")

    if issues:
        for item in issues:
            st.caption(f"• {item}")
    if warnings:
        for item in warnings:
            st.caption(f"• {item}")


def render_confidence_table():
    confidence = st.session_state.get("confidence_map") or {}
    if not confidence:
        return

    st.markdown("#### Confidence")
    rows = []
    for field, meta in confidence.items():
        rows.append({
            "Field": field,
            "Confidence": meta.get("label", "-"),
            "Reason": meta.get("reason", "-")
        })
    if rows:
        st.dataframe(pd.DataFrame(rows), use_container_width=True, height=220, hide_index=True)


def render_invoice_review_form():
    data = st.session_state.get("review_data") or {}
    with st.form("invoice_review_form"):
        c1, c2 = st.columns(2)
        vendor = c1.text_input("Vendor", value=str(data.get("vendor") or data.get("supplier") or ""))
        invoice_number = c2.text_input("Invoice Number", value=str(data.get("invoice_number") or data.get("invoice_no") or ""))
        c3, c4 = st.columns(2)
        invoice_date = c3.text_input("Invoice Date", value=str(data.get("invoice_date") or ""))
        due_date = c4.text_input("Due Date", value=str(data.get("due_date") or ""))
        c5, c6, c7, c8 = st.columns(4)
        currency = c5.text_input("Currency", value=str(data.get("currency") or ""))
        subtotal = c6.text_input("Subtotal", value=str(data.get("subtotal") or ""))
        tax = c7.text_input("Tax", value=str(data.get("tax") or ""))
        total = c8.text_input("Total", value=str(data.get("total") or ""))

        saved = st.form_submit_button("Save Review Changes", use_container_width=True)

    if saved:
        data["vendor"] = vendor
        data["invoice_number"] = invoice_number
        data["invoice_date"] = invoice_date
        data["due_date"] = due_date
        data["currency"] = currency
        data["subtotal"] = subtotal
        data["tax"] = tax
        data["total"] = total
        st.session_state.review_data = data
        refresh_review_scores()
        st.success("Review updates saved")


def render_ticket_review_form():
    data = st.session_state.get("review_data") or {}
    with st.form("ticket_review_form"):
        c1, c2 = st.columns(2)
        traveler_name = c1.text_input("Traveler Name", value=str(data.get("traveler_name") or ""))
        ticket_number = c2.text_input("Ticket Number", value=str(data.get("ticket_number") or ""))
        c3, c4 = st.columns(2)
        airline = c3.text_input("Airline", value=str(data.get("airline") or ""))
        booking_reference = c4.text_input("Booking Reference", value=str(data.get("booking_reference") or ""))
        c5, c6 = st.columns(2)
        from_city = c5.text_input("From", value=str(data.get("from") or ""))
        to_city = c6.text_input("To", value=str(data.get("to") or ""))
        c7, c8, c9 = st.columns(3)
        departure_date = c7.text_input("Departure Date", value=str(data.get("departure_date") or ""))
        return_date = c8.text_input("Return Date", value=str(data.get("return_date") or ""))
        amount = c9.text_input("Amount", value=str(data.get("amount") or ""))

        saved = st.form_submit_button("Save Review Changes", use_container_width=True)

    if saved:
        data["traveler_name"] = traveler_name
        data["ticket_number"] = ticket_number
        data["airline"] = airline
        data["booking_reference"] = booking_reference
        data["from"] = from_city
        data["to"] = to_city
        data["departure_date"] = departure_date
        data["return_date"] = return_date
        data["amount"] = amount
        st.session_state.review_data = data
        refresh_review_scores()
        st.success("Review updates saved")


def render_resume_review_form():
    data = st.session_state.get("review_data") or {}
    with st.form("resume_review_form"):
        c1, c2 = st.columns(2)
        name = c1.text_input("Name", value=str(data.get("name") or ""))
        email = c2.text_input("Email", value=str(data.get("email") or ""))
        c3, c4 = st.columns(2)
        phone = c3.text_input("Phone", value=str(data.get("phone") or ""))
        location = c4.text_input("Location", value=str(data.get("location") or ""))
        linkedin = st.text_input("LinkedIn", value=str(data.get("linkedin") or ""))
        skills = st.text_input("Skills (comma-separated)", value=", ".join(data.get("skills", [])) if isinstance(data.get("skills"), list) else "")
        summary = st.text_area("Summary", value=str(data.get("summary") or ""), height=120)

        saved = st.form_submit_button("Save Review Changes", use_container_width=True)

    if saved:
        data["name"] = name
        data["email"] = email
        data["phone"] = phone
        data["location"] = location
        data["linkedin"] = linkedin
        data["skills"] = [s.strip() for s in skills.split(",") if s.strip()]
        data["summary"] = summary
        st.session_state.review_data = data
        refresh_review_scores()
        st.success("Review updates saved")


def handle_invoice_or_ticket_submission(doc_type):
    validation = st.session_state.get("validation_result") or {}
    if not validation.get("passed"):
        st.warning("Please resolve validation issues before approval")
        return

    data = st.session_state.get("review_data") or {}
    result = send_to_concur(doc_type, data, mode="mock")
    st.session_state.auto_result["result"].update({
        "payload": result.get("payload"),
        "concur_status": result.get("status"),
        "concur_mode": result.get("mode"),
        "concur_submission_id": result.get("submission_id"),
        "concur_batch_id": result.get("batch_id"),
        "concur_document_id": result.get("document_id"),
        "concur_submitted_at": result.get("submitted_at"),
        "concur_endpoint": result.get("endpoint"),
        "concur_processing_state": result.get("processing_state"),
        "concur_next_status": result.get("next_status"),
        "message": result.get("message"),
    })
    st.success(f"{doc_type.title()} approved and submitted to Concur")


def regenerate_resume_from_review():
    validation = st.session_state.get("validation_result") or {}
    data = st.session_state.get("review_data") or {}
    template_bytes = load_default_resume_template_bytes()

    if not template_bytes:
        st.error("Default resume template not found")
        return

    if not validation.get("passed"):
        st.warning("Resume has validation issues. Review before regenerating.")
        return

    try:
        file_bytes = build_resume(data, template_bytes)
        st.session_state.generated_resume = file_bytes
        st.session_state.auto_result["result"]["file"] = file_bytes
        st.session_state.auto_result["result"]["data"] = data
        st.success("Resume regenerated successfully")
    except Exception as e:
        st.error(f"Resume regeneration failed: {str(e)}")

# ------------------------------
# UI RENDERERS
# ------------------------------

def render_header():
    logo_path = Path(__file__).parent / "IDP-Logo1.png"
    col_logo, col_title = st.columns([1, 6], gap="small")

    with col_logo:
        if logo_path.exists():
            st.image(logo_path, width=130)

    with col_title:
        st.markdown("## Intelligent Document Processor")
        st.caption("AI-powered document understanding & automation")


def render_upload_controls():
    with st.sidebar:
        st.markdown("### Account User")
        st.write(f"**{st.session_state['user']}**")
        st.markdown("---")

        st.markdown("### Model")
        model_choice = st.selectbox(
            "Choose Model",
            ["gpt-4o-mini", "gpt-4o", "gpt-5"],
            index=["gpt-4o-mini", "gpt-4o", "gpt-5"].index(
                st.session_state.get("model_choice", "gpt-4o-mini")
            )
        )
        st.session_state["model_choice"] = model_choice

        st.markdown("---")
        st.success("🔑 API key loaded securely")
        
        st.markdown("---")
        cost = st.session_state.get("metrics", {}).get("cost", 0.0)
        st.write(f"💰 Session Cost ${round(cost, 6)}")

        st.markdown("---")

        if st.button("Logout", use_container_width=True):
            for key in ["logged_in", "user", "role", "api_key"]:
                if key in st.session_state:
                    del st.session_state[key]
            st.rerun()

    c1, c2 = st.columns([6, 1], gap="small")
    with c1:
        uploaded_file = st.file_uploader(
            "Upload document",
            type=["txt", "pdf", "docx", "pptx", "xlsx", "png", "jpg", "jpeg"],
            key="main_file_uploader"
        )
    with c2:
        st.write("")
        st.write("")
        if st.button("Reset", use_container_width=True):
            reset_document_state()
            st.rerun()

    st.markdown("---")
    return uploaded_file


def render_agent_activity_panel():
    """
    Kept for compatibility, but not used in the live layout below.
    """
    st.markdown("### Activity")
    st.caption(st.session_state.get("current_step", "Waiting"))
    st.progress(st.session_state.get("progress_value", 0))

    events = st.session_state.get("agent_events", [])
    if not events:
        st.caption("No activity yet")
        return

    for event in events:
        status = event.get("status", "pending")
        icon = "✅" if status == "done" else "❌" if status == "error" else "⏳"
        st.markdown(f"{icon} **{event.get('step', '')}**")
        if event.get("message"):
            st.caption(event["message"])

    with st.expander("Usage", expanded=False):
        auto_metrics = (st.session_state.get("auto_result") or {}).get("metrics", {}) or {}
        c1, c2 = st.columns(2)
        c1.metric("Calls", auto_metrics.get("calls", 0))
        c2.metric("Cost", f"${auto_metrics.get('cost', 0.0):.6f}")

        step_metrics = (st.session_state.get("auto_result") or {}).get("step_metrics", []) or []
        if step_metrics:
            rows = []
            for item in step_metrics:
                rows.append({
                    "Step": item.get("step"),
                    "Cost ($)": round(item.get("metrics", {}).get("cost", 0.0), 6),
                    "Duration (s)": round(item.get("duration", 0.0), 2),
                })
            st.dataframe(pd.DataFrame(rows), use_container_width=True, height=180, hide_index=True)

    with st.expander("Logs", expanded=False):
        for log in st.session_state.get("agent_logs", [])[-20:]:
            st.write(f"- {log}")


def render_empty_state():
    st.markdown("### Ready")
    st.caption("Upload a file to start processing.")


def render_invoice_result(result):
    data = st.session_state.get("review_data") or {}
    st.success(result.get("message", "Invoice processed"))

    c1, c2, c3, c4 = st.columns(4)
    with c1:
        compact_field("Vendor", str(data.get("vendor") or data.get("supplier") or "-"))
    with c2:
        compact_field("Invoice No", str(data.get("invoice_number") or data.get("invoice_no") or "-"))
    with c3:
        compact_field("Date", str(data.get("invoice_date") or "-"))
    with c4:
        compact_field("Total", str(data.get("total") or "-"))

    render_validation_summary()
    render_confidence_table()

    with st.expander("Review & Edit", expanded=True):
        render_invoice_review_form()

    c5, c6 = st.columns([1, 1])
    with c5:
        if st.button("Approve & Send to Concur", use_container_width=True, key="invoice_send"):
            handle_invoice_or_ticket_submission("invoice")
    with c6:
        excel = result.get("excel")
        if excel:
            st.download_button(
                "Download Excel",
                excel,
                f"{(data.get('invoice_number') or data.get('vendor') or 'invoice_data')}.xlsx",
                use_container_width=True
            )

    if result.get("concur_status"):
        st.markdown("#### Concur Delivery")
        c7, c8, c9 = st.columns(3)
        with c7:
            compact_field("Status", str(result.get("concur_status", "-")).title())
        with c8:
            compact_field("Mode", str(result.get("concur_mode", "-")).upper())
        with c9:
            compact_field("Submission ID", str(result.get("concur_submission_id", "-")))

        st.markdown(
            f"<small><b>Batch ID:</b> {result.get('concur_batch_id', '-')} | "
            f"<b>Document ID:</b> {result.get('concur_document_id', '-')}</small>",
            unsafe_allow_html=True
        )
        st.markdown(
            f"<small><b>Submitted At:</b> {result.get('concur_submitted_at', '-')} | "
            f"<b>Processing:</b> {result.get('concur_processing_state', '-')}</small>",
            unsafe_allow_html=True
        )


def render_ticket_result(result):
    data = st.session_state.get("review_data") or {}
    st.success(result.get("message", "Ticket processed"))

    c1, c2, c3, c4 = st.columns(4)
    with c1:
        compact_field("Traveler", str(data.get("traveler_name") or "-"))
    with c2:
        compact_field("Airline", str(data.get("airline") or "-"))
    with c3:
        compact_field("Route", f"{data.get('from', '-')}" + " → " + f"{data.get('to', '-')}")
    with c4:
        compact_field("Amount", str(data.get("amount") or "-"))

    render_validation_summary()
    render_confidence_table()

    with st.expander("Review & Edit", expanded=True):
        render_ticket_review_form()

    if st.button("Approve & Send to Concur", use_container_width=True, key="ticket_send"):
        handle_invoice_or_ticket_submission("ticket")

    if result.get("concur_status"):
        st.markdown("#### Concur Delivery")
        c5, c6, c7 = st.columns(3)
        with c5:
            compact_field("Status", str(result.get("concur_status", "-")).title())
        with c6:
            compact_field("Mode", str(result.get("concur_mode", "-")).upper())
        with c7:
            compact_field("Submission ID", str(result.get("concur_submission_id", "-")))

        st.markdown(
            f"<small><b>Batch ID:</b> {result.get('concur_batch_id', '-')} | "
            f"<b>Document ID:</b> {result.get('concur_document_id', '-')}</small>",
            unsafe_allow_html=True
        )
        st.markdown(
            f"<small><b>Submitted At:</b> {result.get('concur_submitted_at', '-')} | "
            f"<b>Processing:</b> {result.get('concur_processing_state', '-')}</small>",
            unsafe_allow_html=True
        )


def render_resume_result(result):
    data = st.session_state.get("review_data") or {}
    st.success(result.get("message", "Resume generated"))

    top1, top2 = st.columns([2, 1])
    with top1:
        st.caption(f"Output File: {result.get('file_name', 'generated_resume.docx')}")
    with top2:
        if result.get("file"):
            st.download_button(
                "Download Resume",
                data=result["file"],
                file_name=result.get("file_name", "generated_resume.docx"),
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True
            )

    c1, c2 = st.columns([1, 1])
    with c1:
        st.markdown("#### Candidate")
        st.write(data.get("name", ""))
        st.caption(data.get("email", ""))
        st.caption(data.get("phone", ""))
        st.caption(data.get("location", ""))
        st.markdown("#### Skills")
        skills = data.get("skills", [])
        st.caption(", ".join(skills[:15]) if skills else "No skills found")

    with c2:
        st.markdown("#### Summary")
        st.text_area(
            "Summary",
            value=str(data.get("summary", "")),
            height=110,
            label_visibility="collapsed"
        )

    render_validation_summary()
    render_confidence_table()

    with st.expander("Review & Edit", expanded=True):
        render_resume_review_form()

    if st.button("Regenerate Resume", use_container_width=True, key="resume_regen"):
        regenerate_resume_from_review()


def render_generic_result(result):
    st.success(result.get("message", "Processing completed"))
    text = st.session_state.get("full_text", "")
    if text:
        st.text_area("Preview", value=text[:2500], height=180, label_visibility="collapsed")


def render_result_workspace():
    st.markdown("### Result")

    if st.session_state.get("processing_error"):
        st.error(st.session_state["processing_error"])
        return

    if not st.session_state.get("auto_result"):
        render_empty_state()
        return

    doc_type = st.session_state.get("doc_type")
    result = st.session_state.get("auto_result", {}).get("result", {})

    if doc_type == "resume":
        render_resume_result(result)
    elif doc_type == "invoice":
        render_invoice_result(result)
    elif doc_type == "ticket":
        render_ticket_result(result)
    else:
        render_generic_result(result)


def render_chat_section():
    if st.session_state.get("vectorstore") is None:
        st.caption("Upload and process a document first")
        return

    suggested = st.session_state.get("suggested_questions", [])
    if suggested:
        cols = st.columns(len(suggested))
        for i, q in enumerate(suggested):
            if cols[i].button(q, key=f"detail_suggest_{i}"):
                st.session_state.chat_history.append({"role": "user", "content": q})
                docs = st.session_state.vectorstore.similarity_search(
                    q, k=2, filter={"source": st.session_state.current_file}
                )
                context = "\n\n".join([d.page_content[:800] for d in docs])
                response = tracked_llm_call(
                    f"Answer strictly from context.\nContext:\n{context}\nQ:{q}"
                ).content
                st.session_state.chat_history.append({"role": "assistant", "content": response})
                st.rerun()

    for msg in st.session_state.get("chat_history", []):
        with st.chat_message(msg["role"]):
            st.write(msg["content"])

    query = st.chat_input("Ask about this document")
    if query:
        st.session_state.chat_history.append({"role": "user", "content": query})
        docs = st.session_state.vectorstore.similarity_search(
            query, k=2, filter={"source": st.session_state.current_file}
        )
        context = "\n\n".join([d.page_content for d in docs])
        response = tracked_llm_call(
            f"Answer strictly from context.\nContext:\n{context}\nQ:{query}"
        ).content
        st.session_state.chat_history.append({"role": "assistant", "content": response})
        st.rerun()


def render_metrics_section():
    m = st.session_state.get("metrics", {})

    c1, c2, c3, c4 = st.columns(4)
    c1.metric("Total Cost", f"${m.get('cost', 0.0):.6f}")
    c2.metric("Total Tokens", m.get("tokens", 0))
    c3.metric("Input Tokens", m.get("input_tokens", 0))
    c4.metric("Output Tokens", m.get("output_tokens", 0))

    if m.get("response_times"):
        df = pd.DataFrame({
            "Call": list(range(len(m["response_times"]))),
            "Response Time": m["response_times"]
        })
        st.line_chart(df.set_index("Call"))

    st.markdown("#### Cost per Document")
    doc_costs = st.session_state.get("doc_costs", {})
    doc_data = [
        {
            "SL No": i + 1,
            "Document": k,
            "Cost": round(v.get("cost", 0.0), 6),
            "Tokens": v.get("tokens", 0),
        }
        for i, (k, v) in enumerate(doc_costs.items())
    ]

    if doc_data:
        doc_df = pd.DataFrame(doc_data)
        st.dataframe(doc_df, use_container_width=True, hide_index=True, height=220)
        st.bar_chart(doc_df.set_index("Document")[["Cost"]])
    else:
        st.caption("No document-level cost recorded yet")


def render_details_section():
    st.markdown("---")
    st.markdown("### Details")

    with st.expander("Extracted Text", expanded=False):
        full_text = st.session_state.get("full_text")
        if full_text:
            st.text_area("Extracted Text", full_text, height=300, label_visibility="collapsed")

    with st.expander("Structured JSON", expanded=False):
        doc_type = st.session_state.get("doc_type")
        data = st.session_state.get("review_data")
        if doc_type in ["invoice", "ticket", "resume"] and data:
            st.json(data)
        else:
            st.caption("No structured data available")

    with st.expander("Concur Delivery", expanded=False):
        result = (st.session_state.get("auto_result") or {}).get("result", {})
        if st.session_state.get("doc_type") in ["invoice", "ticket"]:
            st.write(f"Status: {result.get('concur_status', 'awaiting approval')}")
            st.write(f"Mode: {result.get('concur_mode', '-')}")
            st.write(f"Submission ID: {result.get('concur_submission_id', '-')}")
            st.write(f"Batch ID: {result.get('concur_batch_id', '-')}")
            st.write(f"Document ID: {result.get('concur_document_id', '-')}")
            if result.get("payload"):
                st.json(result["payload"])
        else:
            st.caption("Concur applies only to invoice and ticket")

    with st.expander("Document Chat", expanded=False):
        render_chat_section()

    with st.expander("Metrics", expanded=False):
        render_metrics_section()


# ------------------------------
# MAIN
# ------------------------------
render_header()
uploaded_file = render_upload_controls()

left_col, right_col = st.columns([1, 1.6], gap="large")

with left_col:
    activity_container = st.container()

    with activity_container:
        st.markdown("### Activity")
        st.session_state["live_step_placeholder"] = st.empty()
        st.session_state["live_progress_placeholder"] = st.empty()
        st.session_state["live_event_placeholder"] = st.empty()

    refresh_live_activity()

    if uploaded_file:
        file_hash = hashlib.md5(uploaded_file.getvalue()).hexdigest()

        if st.session_state.get("file_hash") != file_hash:
            reset_document_state()
            st.session_state.file_hash = file_hash

            with activity_container:
                st.session_state["live_step_placeholder"] = st.empty()
                st.session_state["live_progress_placeholder"] = st.empty()
                st.session_state["live_event_placeholder"] = st.empty()

            refresh_live_activity()

            try:
                process_uploaded_document(uploaded_file)
                refresh_live_activity()
            except Exception as e:
                st.session_state.processing_error = str(e)
                st.session_state.agent_status = "Failed"
                st.session_state.current_step = "Failed"
                st.session_state.elapsed_time = (
                    time.time() - st.session_state.run_started_at
                    if st.session_state.run_started_at else 0.0
                )
                push_agent_log(f"Error: {str(e)}")
                record_agent_event("Processing failed", "error", str(e))
                refresh_live_activity()
                st.error(str(e))
                st.code(traceback.format_exc())
                
with right_col:
    render_result_workspace()

render_details_section()
